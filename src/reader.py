import PyPDF2
from typing import List, Dict
import textract
import os
import docx
import csv
from pptx import Presentation


def read_pdf(path: str) -> List[Dict]:
    result = []
    
    try:
        with open(path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                
                result.append({
                    "page": page_num + 1,
                    "text": text
                })
    
    except FileNotFoundError:
        print(f"Error: The file '{path}' was not found.")
    except PyPDF2.errors.PdfReadError:
        print(f"Error: '{path}' is not a valid PDF file or it's encrypted.")
    except Exception as e:
        print(f"An unexpected error occurred: {str(e)}")
    
    return result


def read_doc(path: str) -> str:
    if not os.path.exists(path):
        raise FileNotFoundError(f"The file '{path}' was not found.")

    file_extension = os.path.splitext(path)[1].lower()

    try:
        if file_extension == '.docx':
            doc = docx.Document(path)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            return '\n'.join(full_text)
        
        elif file_extension == '.doc':
            text = textract.process(path).decode('utf-8')
            return text
        
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")

    except Exception as e:
        raise RuntimeError(f"An error occurred while reading the file: {str(e)}")


def read_csv(path: str) -> str:
    try:
        with open(path, 'r', newline='', encoding='utf-8') as file:
            csv_reader = csv.reader(file)
            header_row = next(csv_reader)
            return ','.join(header_row)
    except StopIteration:
        raise ValueError("The CSV file is empty.")
    except FileNotFoundError:
        raise FileNotFoundError(f"The file '{path}' was not found.")
    except csv.Error as e:
        raise ValueError(f"Error reading CSV file: {str(e)}")
    except Exception as e:
        raise RuntimeError(f"An unexpected error occurred: {str(e)}")


def read_pptx(path: str) -> List[Dict[str, str]]:
    try:
        prs = Presentation(path)
        slides_list = []
        for slide_number, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    slide_text.append(shape.text)
            slides_list.append({
                "slide": slide_number,
                "text": " ".join(slide_text)
            })
        return slides_list
    
    except FileNotFoundError:
        raise FileNotFoundError(f"The file '{path}' was not found.")
    except Exception as e:
        raise RuntimeError(f"An error occurred while reading the PPTX file: {str(e)}")